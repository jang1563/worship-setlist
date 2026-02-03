"""
Export Service for generating various output formats.

Supports:
- ProPresenter 7 format (.pro)
- Plain text lyrics
- PDF (via HTML)
- OpenLyrics XML
- PowerPoint (.pptx)
"""
import json
import re
import io
import xml.etree.ElementTree as ET
from datetime import datetime
from typing import Optional
from dataclasses import dataclass

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


@dataclass
class ExportedSong:
    """A song ready for export."""
    title: str
    artist: str
    key: str
    lyrics: list[dict]  # List of {section, content}
    chords: Optional[str] = None


class ExportService:
    """Service for exporting songs and setlists to various formats."""

    def export_to_propresenter(
        self,
        songs: list[ExportedSong],
        setlist_name: str = "Worship Setlist"
    ) -> str:
        """
        Export songs to ProPresenter 7 JSON format.

        ProPresenter 7 uses a JSON-based format for song data.
        """
        pp_data = {
            "name": setlist_name,
            "created": datetime.now().isoformat(),
            "version": "700",
            "songs": []
        }

        for song in songs:
            pp_song = {
                "title": song.title,
                "artist": song.artist,
                "key": song.key,
                "slides": []
            }

            for section in song.lyrics:
                slide = {
                    "label": section.get("section", ""),
                    "text": section.get("content", ""),
                    "notes": ""
                }
                pp_song["slides"].append(slide)

            pp_data["songs"].append(pp_song)

        return json.dumps(pp_data, ensure_ascii=False, indent=2)

    def export_to_openlyrics(self, song: ExportedSong) -> str:
        """
        Export a song to OpenLyrics XML format.

        OpenLyrics is an open standard for song lyrics interchange.
        """
        # Create root element
        root = ET.Element("song")
        root.set("xmlns", "http://openlyrics.info/namespace/2009/song")
        root.set("version", "0.9")
        root.set("createdIn", "송플래너")
        root.set("modifiedDate", datetime.now().strftime("%Y-%m-%dT%H:%M:%S"))

        # Properties
        properties = ET.SubElement(root, "properties")

        titles = ET.SubElement(properties, "titles")
        title = ET.SubElement(titles, "title")
        title.text = song.title

        if song.artist:
            authors = ET.SubElement(properties, "authors")
            author = ET.SubElement(authors, "author")
            author.text = song.artist

        if song.key:
            key_elem = ET.SubElement(properties, "key")
            key_elem.text = song.key

        # Lyrics
        lyrics_elem = ET.SubElement(root, "lyrics")

        for i, section in enumerate(song.lyrics):
            verse = ET.SubElement(lyrics_elem, "verse")
            verse.set("name", section.get("section", f"v{i+1}"))

            lines = ET.SubElement(verse, "lines")
            for line in section.get("content", "").split("\n"):
                line_elem = ET.SubElement(lines, "line")
                line_elem.text = line.strip()

        # Convert to string
        return ET.tostring(root, encoding="unicode", xml_declaration=True)

    def export_to_plain_text(
        self,
        songs: list[ExportedSong],
        include_chords: bool = False,
        include_headers: bool = True
    ) -> str:
        """Export songs to plain text format."""
        lines = []

        for song in songs:
            if include_headers:
                lines.append("=" * 50)
                lines.append(f"제목: {song.title}")
                lines.append(f"아티스트: {song.artist}")
                lines.append(f"키: {song.key}")
                lines.append("=" * 50)
                lines.append("")

            for section in song.lyrics:
                if section.get("section"):
                    lines.append(f"[{section['section']}]")

                content = section.get("content", "")
                if not include_chords:
                    # Remove chord brackets
                    content = re.sub(r'\[[^\]]+\]', '', content)

                lines.append(content)
                lines.append("")

            lines.append("\n")

        return "\n".join(lines)

    def export_setlist_to_html(
        self,
        songs: list[ExportedSong],
        setlist_name: str,
        date: Optional[str] = None,
        service_type: Optional[str] = None
    ) -> str:
        """Export setlist to printable HTML format."""
        html = f"""<!DOCTYPE html>
<html lang="ko">
<head>
    <meta charset="UTF-8">
    <title>{setlist_name}</title>
    <style>
        body {{
            font-family: 'Pretendard', -apple-system, sans-serif;
            max-width: 800px;
            margin: 0 auto;
            padding: 20px;
            color: #1f2937;
        }}
        .header {{
            text-align: center;
            border-bottom: 2px solid #333;
            padding-bottom: 20px;
            margin-bottom: 30px;
        }}
        .header h1 {{
            margin: 0;
            font-size: 24px;
        }}
        .header .meta {{
            color: #6b7280;
            margin-top: 8px;
        }}
        .song {{
            page-break-inside: avoid;
            margin-bottom: 40px;
        }}
        .song-header {{
            display: flex;
            justify-content: space-between;
            align-items: center;
            border-bottom: 1px solid #e5e7eb;
            padding-bottom: 8px;
            margin-bottom: 16px;
        }}
        .song-title {{
            font-size: 18px;
            font-weight: 600;
        }}
        .song-key {{
            background: #f3f4f6;
            padding: 4px 12px;
            border-radius: 4px;
            font-family: monospace;
            font-weight: 600;
        }}
        .section {{
            margin-bottom: 16px;
        }}
        .section-label {{
            font-size: 12px;
            text-transform: uppercase;
            color: #6b7280;
            margin-bottom: 4px;
        }}
        .lyrics {{
            white-space: pre-wrap;
            line-height: 1.8;
        }}
        .chord {{
            color: #4f46e5;
            font-weight: 600;
            font-family: monospace;
        }}
        .footer {{
            text-align: center;
            margin-top: 40px;
            padding-top: 20px;
            border-top: 1px solid #e5e7eb;
            font-size: 12px;
            color: #9ca3af;
        }}
        @media print {{
            body {{
                padding: 0;
            }}
            .song {{
                page-break-after: always;
            }}
            .song:last-child {{
                page-break-after: auto;
            }}
        }}
    </style>
</head>
<body>
    <div class="header">
        <h1>{setlist_name}</h1>
        <div class="meta">
            {f'{date} · ' if date else ''}{service_type or '예배'}
        </div>
    </div>
"""

        for i, song in enumerate(songs):
            html += f"""
    <div class="song">
        <div class="song-header">
            <span class="song-title">{i+1}. {song.title}</span>
            <span class="song-key">{song.key}</span>
        </div>
"""
            for section in song.lyrics:
                section_label = section.get("section", "")
                content = self._format_chords_html(section.get("content", ""))
                html += f"""
        <div class="section">
            <div class="section-label">{section_label}</div>
            <div class="lyrics">{content}</div>
        </div>
"""
            html += "    </div>\n"

        html += f"""
    <div class="footer">
        송플래너 | 생성일: {datetime.now().strftime('%Y-%m-%d')}
    </div>
</body>
</html>
"""
        return html

    def _format_chords_html(self, text: str) -> str:
        """Format chord brackets as styled spans."""
        return re.sub(
            r'\[([^\]]+)\]',
            r'<span class="chord">[\1]</span>',
            text
        )

    def chordpro_to_sections(self, chordpro: str) -> list[dict]:
        """Convert ChordPro content to section list."""
        sections = []
        current_section = "Verse"
        current_lines = []

        for line in chordpro.split('\n'):
            # Check for comment (section marker)
            comment_match = re.match(r'\{comment:\s*(.+?)\}', line, re.IGNORECASE)
            if comment_match:
                # Save previous section
                if current_lines:
                    sections.append({
                        "section": current_section,
                        "content": '\n'.join(current_lines)
                    })
                current_section = comment_match.group(1)
                current_lines = []
                continue

            # Skip other directives
            if line.strip().startswith('{'):
                continue

            # Add to current section
            if line.strip():
                current_lines.append(line)

        # Save last section
        if current_lines:
            sections.append({
                "section": current_section,
                "content": '\n'.join(current_lines)
            })

        return sections

    def export_song_to_pdf_html(
        self,
        song: ExportedSong,
        include_chords: bool = True
    ) -> str:
        """Export a single song to PDF-optimized HTML format."""
        html = f"""<!DOCTYPE html>
<html lang="ko">
<head>
    <meta charset="UTF-8">
    <title>{song.title}</title>
    <style>
        @page {{
            size: A4;
            margin: 20mm;
        }}
        body {{
            font-family: 'Pretendard', 'Noto Sans KR', -apple-system, sans-serif;
            font-size: 14px;
            line-height: 1.6;
            color: #000;
            max-width: 210mm;
            margin: 0 auto;
        }}
        .song-header {{
            text-align: center;
            margin-bottom: 30px;
            padding-bottom: 15px;
            border-bottom: 2px solid #333;
        }}
        .song-title {{
            font-size: 28px;
            font-weight: 700;
            margin: 0 0 8px 0;
        }}
        .song-artist {{
            font-size: 16px;
            color: #555;
            margin: 0 0 8px 0;
        }}
        .song-key {{
            display: inline-block;
            background: #f0f0f0;
            padding: 4px 16px;
            border-radius: 4px;
            font-family: monospace;
            font-weight: 600;
            font-size: 16px;
        }}
        .section {{
            margin-bottom: 24px;
        }}
        .section-label {{
            font-size: 12px;
            font-weight: 600;
            text-transform: uppercase;
            color: #666;
            margin-bottom: 8px;
            padding: 2px 8px;
            background: #f5f5f5;
            display: inline-block;
        }}
        .lyrics-container {{
            padding-left: 16px;
        }}
        .lyric-line {{
            margin-bottom: 4px;
        }}
        .chord-line {{
            font-family: 'Consolas', 'Monaco', monospace;
            color: #0066cc;
            font-weight: 600;
            font-size: 13px;
            margin-bottom: 2px;
        }}
        .text-line {{
            font-size: 16px;
        }}
        .footer {{
            margin-top: 40px;
            padding-top: 15px;
            border-top: 1px solid #ddd;
            text-align: center;
            font-size: 10px;
            color: #999;
        }}
    </style>
</head>
<body>
    <div class="song-header">
        <h1 class="song-title">{song.title}</h1>
        <p class="song-artist">{song.artist}</p>
        <span class="song-key">Key: {song.key}</span>
    </div>
"""
        for section in song.lyrics:
            section_label = section.get("section", "")
            content = section.get("content", "")

            html += f"""    <div class="section">
        <div class="section-label">{section_label}</div>
        <div class="lyrics-container">
"""
            if include_chords:
                # Parse chords and lyrics separately
                lines = content.split('\n')
                for line in lines:
                    chord_text, lyric_text = self._split_chord_lyric_line(line)
                    if chord_text:
                        html += f'            <div class="chord-line">{chord_text}</div>\n'
                    html += f'            <div class="text-line">{lyric_text}</div>\n'
            else:
                # Remove chords
                clean_content = re.sub(r'\[[^\]]+\]', '', content)
                for line in clean_content.split('\n'):
                    html += f'            <div class="text-line">{line}</div>\n'

            html += """        </div>
    </div>
"""

        html += f"""    <div class="footer">
        WorshipFlow 찬양설계 | {datetime.now().strftime('%Y-%m-%d')}
    </div>
</body>
</html>
"""
        return html

    def _split_chord_lyric_line(self, line: str) -> tuple[str, str]:
        """Split a ChordPro line into chord line and lyric line."""
        if '[' not in line:
            return '', line

        chord_positions = []
        lyric_parts = []
        current_pos = 0

        for match in re.finditer(r'\[([^\]]+)\]', line):
            # Add lyrics before this chord
            lyric_parts.append(line[current_pos:match.start()])
            # Record chord position
            lyric_len = len(''.join(lyric_parts))
            chord_positions.append((lyric_len, match.group(1)))
            current_pos = match.end()

        # Add remaining lyrics
        lyric_parts.append(line[current_pos:])
        lyric_text = ''.join(lyric_parts)

        # Build chord line with proper spacing
        chord_line = []
        last_pos = 0
        for pos, chord in chord_positions:
            spaces_needed = pos - last_pos
            chord_line.append(' ' * max(0, spaces_needed))
            chord_line.append(chord)
            last_pos = pos + len(chord)

        return ''.join(chord_line), lyric_text

    def export_setlist_summary_html(
        self,
        songs: list[dict],
        setlist_name: str,
        date: Optional[str] = None,
        service_type: Optional[str] = None,
        total_duration_min: Optional[int] = None
    ) -> str:
        """Export a setlist summary (song list only) for printing."""
        html = f"""<!DOCTYPE html>
<html lang="ko">
<head>
    <meta charset="UTF-8">
    <title>{setlist_name} - 요약</title>
    <style>
        @page {{ size: A4; margin: 15mm; }}
        body {{
            font-family: 'Pretendard', 'Noto Sans KR', sans-serif;
            font-size: 12px;
            max-width: 210mm;
            margin: 0 auto;
            padding: 20px;
        }}
        .header {{
            text-align: center;
            margin-bottom: 30px;
        }}
        .header h1 {{ font-size: 22px; margin: 0 0 8px 0; }}
        .header .meta {{ color: #666; }}
        table {{
            width: 100%;
            border-collapse: collapse;
            margin-top: 20px;
        }}
        th, td {{
            padding: 10px 8px;
            text-align: left;
            border-bottom: 1px solid #ddd;
        }}
        th {{
            background: #f5f5f5;
            font-weight: 600;
        }}
        .song-num {{ width: 40px; text-align: center; }}
        .song-key {{ font-family: monospace; font-weight: 600; }}
        .song-duration {{ text-align: right; }}
        .total-row {{ font-weight: 600; background: #f0f0f0; }}
    </style>
</head>
<body>
    <div class="header">
        <h1>{setlist_name}</h1>
        <div class="meta">{date or datetime.now().strftime('%Y-%m-%d')} · {service_type or '예배'}</div>
    </div>
    <table>
        <thead>
            <tr>
                <th class="song-num">#</th>
                <th>곡 제목</th>
                <th>아티스트</th>
                <th>키</th>
                <th>역할</th>
                <th class="song-duration">시간</th>
            </tr>
        </thead>
        <tbody>
"""
        total_seconds = 0
        for i, song in enumerate(songs):
            duration_sec = song.get('duration_sec', 0)
            total_seconds += duration_sec
            duration_str = f"{duration_sec // 60}:{(duration_sec % 60):02d}" if duration_sec else "-"

            html += f"""            <tr>
                <td class="song-num">{i + 1}</td>
                <td>{song.get('title', '')}</td>
                <td>{song.get('artist', '')}</td>
                <td class="song-key">{song.get('key', '')}</td>
                <td>{song.get('role', '')}</td>
                <td class="song-duration">{duration_str}</td>
            </tr>
"""

        total_min = total_duration_min or (total_seconds // 60)
        html += f"""        </tbody>
        <tfoot>
            <tr class="total-row">
                <td colspan="5">총 {len(songs)}곡</td>
                <td class="song-duration">{total_min}분</td>
            </tr>
        </tfoot>
    </table>
</body>
</html>
"""
        return html


    def export_to_powerpoint(
        self,
        songs: list[ExportedSong],
        setlist_name: str = "예배 찬양",
        include_chords: bool = False
    ) -> Optional[bytes]:
        """Export songs to PowerPoint format.

        Returns PowerPoint file as bytes, or None if pptx not available.
        """
        if not PPTX_AVAILABLE:
            return None

        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)

        # Title slide
        title_slide_layout = prs.slide_layouts[6]  # Blank slide
        slide = prs.slides.add_slide(title_slide_layout)

        # Add title text
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3), Inches(15), Inches(2)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = setlist_name
        title_para.font.size = Pt(60)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

        # Add date
        date_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5), Inches(15), Inches(1)
        )
        date_frame = date_box.text_frame
        date_para = date_frame.paragraphs[0]
        date_para.text = datetime.now().strftime('%Y년 %m월 %d일')
        date_para.font.size = Pt(28)
        date_para.font.color.rgb = RGBColor(200, 200, 200)
        date_para.alignment = PP_ALIGN.CENTER

        # Set background color (dark)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(30, 30, 50)

        # Song slides
        for song in songs:
            # Song title slide
            slide = prs.slides.add_slide(title_slide_layout)
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(30, 30, 50)

            # Song title
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(3.5), Inches(15), Inches(1.5)
            )
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = song.title
            title_para.font.size = Pt(54)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255)
            title_para.alignment = PP_ALIGN.CENTER

            # Artist and key
            info_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(5), Inches(15), Inches(1)
            )
            info_frame = info_box.text_frame
            info_para = info_frame.paragraphs[0]
            info_para.text = f"{song.artist}  |  Key: {song.key}"
            info_para.font.size = Pt(24)
            info_para.font.color.rgb = RGBColor(180, 180, 180)
            info_para.alignment = PP_ALIGN.CENTER

            # Lyrics slides
            for section in song.lyrics:
                content = section.get("content", "")
                if not include_chords:
                    content = re.sub(r'\[[^\]]+\]', '', content)

                # Split into chunks for slides (max 6 lines per slide)
                lines = content.strip().split('\n')
                chunks = [lines[i:i+6] for i in range(0, len(lines), 6)]

                for chunk in chunks:
                    slide = prs.slides.add_slide(title_slide_layout)
                    fill = slide.background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(0, 0, 0)

                    # Section label (small, top)
                    if section.get("section"):
                        label_box = slide.shapes.add_textbox(
                            Inches(0.5), Inches(0.3), Inches(15), Inches(0.5)
                        )
                        label_frame = label_box.text_frame
                        label_para = label_frame.paragraphs[0]
                        label_para.text = section["section"]
                        label_para.font.size = Pt(16)
                        label_para.font.color.rgb = RGBColor(150, 150, 150)
                        label_para.alignment = PP_ALIGN.CENTER

                    # Lyrics content
                    lyrics_text = '\n'.join(chunk)
                    lyrics_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(2), Inches(15), Inches(6)
                    )
                    lyrics_frame = lyrics_box.text_frame
                    lyrics_frame.word_wrap = True

                    lyrics_para = lyrics_frame.paragraphs[0]
                    lyrics_para.text = lyrics_text
                    lyrics_para.font.size = Pt(48)
                    lyrics_para.font.color.rgb = RGBColor(255, 255, 255)
                    lyrics_para.alignment = PP_ALIGN.CENTER
                    lyrics_para.line_spacing = 1.5

        # Save to bytes
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        return pptx_buffer.getvalue()


# Singleton instance
export_service = ExportService()
